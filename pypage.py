from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_swo_presentation():
    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Color scheme (Nord theme)
    colors = {
        'dark_bg': RGBColor(46, 52, 64),     # #2e3440
        'light_text': RGBColor(216, 222, 233), # #d8dee9
        'blue': RGBColor(136, 192, 208),      # #88c0d0
    }

    # --- Helper Function for Safe Image Addition ---
    def add_safe_image(slide, img_path, left, top, width, height):
        """Adds image with error handling"""
        try:
            if os.path.exists(img_path):
                slide.shapes.add_picture(img_path, left, top, width, height)
            else:
                # Create placeholder if image missing
                txt_box = slide.shapes.add_textbox(left, top, width, height)
                tf = txt_box.text_frame
                tf.text = f"Image not found:\n{os.path.basename(img_path)}"
                tf.paragraphs[0].font.color.rgb = colors['light_text']
                print(f"Warning: Missing image - {img_path}")
        except Exception as e:
            print(f"Error adding image: {e}")

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = colors['dark_bg']

    title = slide.shapes.title
    title.text = "SWO System Documentation"
    title.text_frame.paragraphs[0].font.color.rgb = colors['light_text']
    title.text_frame.paragraphs[0].font.size = Pt(44)

    # --- Slide 2: System Mental Model ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = colors['dark_bg']

    title = slide.shapes.title
    title.text = "System Mental Model"
    title.text_frame.paragraphs[0].font.color.rgb = colors['blue']

    # Add image (with error handling)
    img_path = os.path.join("images", "mindmap.png")  # Relative path
    add_safe_image(
        slide,
        img_path=img_path,
        left=Inches(1),
        top=Inches(1.5),
        width=Inches(8),
        height=Inches(4)
    )

    # --- Slide 3: SOP Steps ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = colors['dark_bg']

    title = slide.shapes.title
    title.text = "Monthly SWO Update SOP"
    title.text_frame.paragraphs[0].font.color.rgb = colors['blue']

    steps = [
        "1. Data Collection from Outlook",
        "2. Data Validation",
        "3. Schema Compliance Check",
        "4. Data Transformation",
        "5. Database Update",
        "6. Verification",
        "7. Documentation"
    ]

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    tf = textbox.text_frame
    for step in steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.color.rgb = colors['light_text']
        p.font.size = Pt(18)
        p.space_after = Pt(12)

    # Save presentation
    output_path = "SWO_Documentation.pptx"
    try:
        prs.save(output_path)
        print(f"Successfully created: {output_path}")
    except Exception as e:
        print(f"Error saving presentation: {e}")

if __name__ == "__main__":
    # Create images directory if it doesn't exist
    if not os.path.exists("images"):
        os.makedirs("images")
        print("Created 'images' directory for diagrams")
    
    create_swo_presentation()